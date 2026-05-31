# -*- coding: utf-8 -*-
"""100회 주기(같은 끝 두 자리 회차) 당첨번호 재출현 통계 → PPT.
대상 회차 1127~1226 각각의 당첨번호가 100회 이전 형제 회차들에 나왔는지 집계."""
import csv, sys
from collections import Counter
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

TARGET_LO, TARGET_HI = 1127, 1226
# 조회 범위(회). 인자로 변경: python cycle_ppt.py 200  → 200회 이내(T-100·T-200)
MAX_LOOKBACK = int(sys.argv[1]) if len(sys.argv) > 1 else 300
NSIB = MAX_LOOKBACK // 100   # 형제 회차 개수
CSV = "src/resources/lotto.csv"

FONT = "Apple SD Gothic Neo"
NAVY = RGBColor(0x1F, 0x29, 0x4E)
BLUE = RGBColor(0x2E, 0x5C, 0xB8)
GOLD = RGBColor(0xD4, 0xA0, 0x17)
GRAY = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xEE, 0xF2, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x2E, 0x7D, 0x32)

# ---------------- 데이터 & 통계 ----------------
def load(path):
    draws = {}
    with open(path, encoding="utf-8") as f:
        r = csv.reader(f)
        next(r)
        for row in r:
            if len(row) < 7:
                continue
            try:
                d = int(row[0]); nums = [int(x) for x in row[1:7]]
            except ValueError:
                continue
            draws[d] = sorted(nums)
    return draws

def siblings(T):
    s = []; d = T - 100
    while d >= 1 and (T - d) <= MAX_LOOKBACK:   # 300회 이내만
        s.append(d); d -= 100
    return sorted(s)

draws = load(CSV)
results = []   # (T, sibs, [(num, [hit draws])...], match_count)
tot_num = tot_with = tot_pair = 0
exp_with = exp_pair = 0.0
for T in range(TARGET_LO, TARGET_HI + 1):
    if T not in draws:
        continue
    sibs = siblings(T); k = len(sibs)
    per = []
    for n in draws[T]:
        hit = [s for s in sibs if n in draws.get(s, [])]
        per.append((n, hit))
        tot_num += 1
        if hit: tot_with += 1
        tot_pair += len(hit)
        exp_pair += k * 6 / 45
        exp_with += 1 - (1 - 6/45) ** k
    results.append((T, sibs, per, sum(1 for _, h in per if h)))

DRAWN = len(results)
dist = Counter(m for _, _, _, m in results)

# 동적 라벨
SIB_DESC = "·".join(f"T−{100*i}" for i in range(1, NSIB + 1))          # 예: T−100·T−200
EX1226 = "·".join(f"{x}" for x in siblings(1226))                       # 예: 1026·1126
OBS_PCT = 100 * tot_with / tot_num
EXP_PCT = 100 * exp_with / tot_num

# ---------------- PPT 헬퍼 ----------------
prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def add_slide(): return prs.slides.add_slide(BLANK)
def set_bg(s, c):
    s.background.fill.solid(); s.background.fill.fore_color.rgb = c

def textbox(s, l, t, w, h, text, size, color=NAVY, bold=False,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(l, t, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; r = p.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color; r.font.name = FONT
    return tb

def bar(s, l, t, w, h, c):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = c
    sp.line.fill.background(); sp.shadow.inherit = False
    return sp

def header(s, title):
    bar(s, 0, 0, SW, Inches(1.0), NAVY)
    bar(s, 0, Inches(1.0), SW, Pt(4), GOLD)
    textbox(s, Inches(0.5), 0, Inches(12.3), Inches(1.0), title,
            24, WHITE, True, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)

def style_cell(cell, size, color, bold=False, align=PP_ALIGN.CENTER, fill=None):
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_top = Pt(1); cell.margin_bottom = Pt(1)
    cell.margin_left = Pt(4); cell.margin_right = Pt(4)
    if fill is not None:
        cell.fill.solid(); cell.fill.fore_color.rgb = fill
    else:
        cell.fill.background()
    for p in cell.text_frame.paragraphs:
        p.alignment = align
        for r in p.runs:
            r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = color; r.font.name = FONT

# ---------------- Slide 1: 표지 ----------------
s = add_slide(); set_bg(s, NAVY)
bar(s, 0, Inches(3.05), SW, Pt(3), GOLD)
textbox(s, Inches(1), Inches(2.0), Inches(11.3), Inches(1.0),
        "100회 주기 당첨번호 재출현 통계", 40, WHITE, True, PP_ALIGN.CENTER)
textbox(s, Inches(1), Inches(3.2), Inches(11.3), Inches(0.7),
        f"대상 회차 {TARGET_LO}~{TARGET_HI}회 ({DRAWN}회) · {MAX_LOOKBACK}회 이내", 22, GOLD, True, PP_ALIGN.CENTER)
textbox(s, Inches(1), Inches(4.15), Inches(11.3), Inches(0.7),
        f"각 회차 당첨번호가 최근 {NSIB}개 형제 회차({SIB_DESC})에 나왔는지 분석",
        16, RGBColor(0xC9, 0xD3, 0xEA), False, PP_ALIGN.CENTER)
textbox(s, Inches(1), Inches(6.6), Inches(11.3), Inches(0.5),
        "2026-05-31 생성", 13, RGBColor(0x9A, 0xA8, 0xC9), False, PP_ALIGN.CENTER)

# ---------------- Slide 2: 방법 + 요약 ----------------
s = add_slide(); set_bg(s, WHITE)
header(s, "분석 방법 및 요약 통계")
textbox(s, Inches(0.6), Inches(1.2), Inches(12.1), Inches(1.0),
        f"방법: 대상 회차 T({TARGET_LO}~{TARGET_HI})의 6개 당첨번호 각각이, {MAX_LOOKBACK}회 이내 형제 회차\n"
        f"({SIB_DESC}, 예: 1226회 → {EX1226}회)의 당첨번호에 있었는지 확인합니다.",
        14, GRAY, False)
rows = [
    ("항목", "관측값", "무작위 기대값"),
    ("대상 회차 수", f"{DRAWN}회", "—"),
    ("총 당첨번호 수", f"{tot_num}개", "—"),
    ("주기 재출현 번호 (형제 1개↑)", f"{tot_with}개 ({100*tot_with/tot_num:.1f}%)", f"{exp_with:.0f}개 ({100*exp_with/tot_num:.1f}%)"),
    ("번호–회차 매칭 총수", f"{tot_pair}개", f"{exp_pair:.0f}개"),
    ("회차당 평균 재출현 번호", f"{tot_with/DRAWN:.2f}개 / 6", f"{exp_with/DRAWN:.2f}개 / 6"),
]
tbl = s.shapes.add_table(len(rows), 3, Inches(0.6), Inches(2.5),
                         Inches(11.0), Inches(0.55)*len(rows)).table
tbl.columns[0].width = Inches(5.2); tbl.columns[1].width = Inches(3.0); tbl.columns[2].width = Inches(2.8)
for j, h in enumerate(rows[0]):
    tbl.cell(0, j).text = h
    style_cell(tbl.cell(0, j), 13, WHITE, True, PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER, fill=BLUE)
for i in range(1, len(rows)):
    fill = LIGHT if i % 2 == 0 else WHITE
    for j, v in enumerate(rows[i]):
        tbl.cell(i, j).text = v
        style_cell(tbl.cell(i, j), 12, NAVY, j == 1, PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER, fill=fill)
obs_pct = 100*tot_with/tot_num
exp_pct = 100*exp_with/tot_num
textbox(s, Inches(0.6), Inches(6.0), Inches(12.1), Inches(1.3),
        f"▶ 결론: 관측 재출현율({obs_pct:.1f}%)이 무작위 기대치({exp_pct:.1f}%)와 거의 같습니다. 즉 '100회 주기 재출현'은\n"
        "   통계적 우연으로 설명되며, 다음 회차 예측에 쓸 수 있는 신호로 보기는 어렵습니다. (참고용 통계)",
        13, RED, True, PP_ALIGN.LEFT)

# ---------------- Slide 3: 회차별 재출현 개수 분포 ----------------
s = add_slide(); set_bg(s, WHITE)
header(s, "회차별 재출현 번호 개수 분포 (6개 중)")
maxv = max(dist.get(i, 0) for i in range(7))
top = Inches(1.55); rowh = Inches(0.72)
for i in range(7):
    cnt = dist.get(i, 0)
    y = top + rowh * i
    textbox(s, Inches(0.6), y, Inches(2.2), rowh, f"{i}개 재출현", 15, NAVY, True,
            PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    w = int(Inches(8.0) * (cnt / maxv)) if maxv else Emu(1)
    if w < Emu(1): w = Emu(1)
    c = GOLD if cnt == maxv else BLUE
    bar(s, Inches(2.9), y + Inches(0.1), w, Inches(0.5), c)
    textbox(s, Inches(2.9) + w + Inches(0.1), y, Inches(2.2), rowh,
            f"{cnt}회 ({100*cnt/DRAWN:.0f}%)", 13, GRAY, True, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
textbox(s, Inches(0.6), Inches(6.75), Inches(12), Inches(0.5),
        f"대상 {DRAWN}회 중 6개 모두 재출현한 회차 {dist.get(6,0)}회 · 5개 {dist.get(5,0)}회 · 4개 {dist.get(4,0)}회",
        13, GRAY, True, PP_ALIGN.LEFT)

# ---------------- Slide 4: 예시 1226회 상세 ----------------
s = add_slide(); set_bg(s, WHITE)
ex = next((r for r in results if r[0] == 1226), None)
if ex:
    T, sibs, per, mc = ex
    header(s, f"예시 — {T}회 당첨번호의 100회 주기 재출현")
    textbox(s, Inches(0.6), Inches(1.2), Inches(12.1), Inches(0.6),
            f"형제 회차: {', '.join(str(x) for x in sibs)}회  (총 {len(sibs)}개)",
            13, GRAY, False)
    erows = [("당첨번호", "재출현한 형제 회차", "횟수")]
    for n, hit in per:
        erows.append((f"{n}번", ", ".join(f"{x}회" for x in hit) if hit else "(없음)", str(len(hit))))
    tbl = s.shapes.add_table(len(erows), 3, Inches(0.6), Inches(2.0),
                             Inches(11.0), Inches(0.6)*len(erows)).table
    tbl.columns[0].width = Inches(2.0); tbl.columns[1].width = Inches(7.5); tbl.columns[2].width = Inches(1.5)
    for j, h in enumerate(erows[0]):
        tbl.cell(0, j).text = h
        style_cell(tbl.cell(0, j), 13, WHITE, True, PP_ALIGN.LEFT if j == 1 else PP_ALIGN.CENTER, fill=BLUE)
    for i in range(1, len(erows)):
        hit_none = erows[i][1] == "(없음)"
        fill = LIGHT if i % 2 == 0 else WHITE
        tbl.cell(i, 0).text = erows[i][0]
        style_cell(tbl.cell(i, 0), 14, NAVY, True, fill=GOLD)
        tbl.cell(i, 1).text = erows[i][1]
        style_cell(tbl.cell(i, 1), 12, GRAY if hit_none else NAVY, False, PP_ALIGN.LEFT, fill=fill)
        tbl.cell(i, 2).text = erows[i][2]
        style_cell(tbl.cell(i, 2), 12, RED, True, fill=fill)

# ---------------- Slide 5+: 전체 100회 상세 ----------------
def detail_str(per):
    parts = []
    for n, hit in per:
        if hit:
            parts.append(f"{n}→{'·'.join(str(x) for x in hit)}")
    return " , ".join(parts) if parts else "(재출현 없음)"

PER_SLIDE = 25
rows_all = sorted(results, key=lambda r: -r[0])  # 최신 회차 먼저
pages = [rows_all[i:i+PER_SLIDE] for i in range(0, len(rows_all), PER_SLIDE)]
for pi, page in enumerate(pages):
    s = add_slide(); set_bg(s, WHITE)
    header(s, f"전체 상세 — 회차별 재출현 번호→회차 ({pi+1}/{len(pages)})")
    hdr = ["대상회차", "재출현", "재출현 번호 → 나온 형제 회차"]
    tbl = s.shapes.add_table(len(page) + 1, 3, Inches(0.25), Inches(1.18),
                             Inches(12.85), Inches(6.05)).table
    tbl.columns[0].width = Inches(1.25); tbl.columns[1].width = Inches(0.95); tbl.columns[2].width = Inches(10.65)
    for j, h in enumerate(hdr):
        tbl.cell(0, j).text = h
        style_cell(tbl.cell(0, j), 11, WHITE, True, PP_ALIGN.LEFT if j == 2 else PP_ALIGN.CENTER, fill=BLUE)
    for i, (T, sibs, per, mc) in enumerate(page, start=1):
        fill = LIGHT if i % 2 == 0 else WHITE
        tbl.cell(i, 0).text = f"{T}회"
        style_cell(tbl.cell(i, 0), 10, NAVY, True, fill=fill)
        tbl.cell(i, 1).text = f"{mc}개"
        style_cell(tbl.cell(i, 1), 10, RED if mc >= 5 else NAVY, True, fill=fill)
        tbl.cell(i, 2).text = detail_str(per)
        style_cell(tbl.cell(i, 2), 8, GRAY, False, PP_ALIGN.LEFT, fill)

out = f"로또_100회주기_{MAX_LOOKBACK}회이내_재출현통계_{TARGET_LO}-{TARGET_HI}.pptx"
prs.save(out)
print("SAVED:", out)
print("slides:", len(prs.slides._sldIdLst))
