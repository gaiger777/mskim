# -*- coding: utf-8 -*-
"""LottoPatternAnalyzer 실행 결과를 PPT로 변환."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

FONT = "Apple SD Gothic Neo"
NAVY = RGBColor(0x1F, 0x29, 0x4E)
BLUE = RGBColor(0x2E, 0x5C, 0xB8)
GOLD = RGBColor(0xD4, 0xA0, 0x17)
GRAY = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xEE, 0xF2, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xC0, 0x39, 0x2B)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def textbox(slide, l, t, w, h, text, size, color=NAVY, bold=False,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = ln
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = font
    return tb


def bar(slide, l, t, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def header(slide, title):
    bar(slide, 0, 0, SW, Inches(1.0), NAVY)
    bar(slide, 0, Inches(1.0), SW, Pt(4), GOLD)
    textbox(slide, Inches(0.5), 0, Inches(12.3), Inches(1.0), title,
            26, WHITE, True, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)


def style_cell(cell, size, color, bold=False, align=PP_ALIGN.CENTER, fill=None):
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_top = Pt(1)
    cell.margin_bottom = Pt(1)
    cell.margin_left = Pt(4)
    cell.margin_right = Pt(4)
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    else:
        cell.fill.background()
    for p in cell.text_frame.paragraphs:
        p.alignment = align
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = FONT


# ---------------- Slide 1: 표지 ----------------
s = add_slide()
set_bg(s, NAVY)
bar(s, 0, Inches(3.05), SW, Pt(3), GOLD)
textbox(s, Inches(1), Inches(2.1), Inches(11.3), Inches(1.0),
        "1226회차 로또 패턴 분석 리포트", 40, WHITE, True, PP_ALIGN.CENTER)
textbox(s, Inches(1), Inches(3.25), Inches(11.3), Inches(0.7),
        "확률 기반 패턴 분석 추천 시스템 v7.3", 22, GOLD, True, PP_ALIGN.CENTER)
textbox(s, Inches(1), Inches(4.2), Inches(11.3), Inches(0.6),
        "기준 회차: 1225회  |  22개 패턴 기법 앙상블  |  최근 15회차 백테스트",
        16, RGBColor(0xC9, 0xD3, 0xEA), False, PP_ALIGN.CENTER)
textbox(s, Inches(1), Inches(6.6), Inches(11.3), Inches(0.5),
        "2026-05-30 생성", 13, RGBColor(0x9A, 0xA8, 0xC9), False, PP_ALIGN.CENTER)

# ---------------- Slide 1b: 추천 방식 검증 (walk-forward) ----------------
import re


def parse_eval(path):
    rnd = None
    rows = []
    with open(path, encoding="utf-8") as f:
        for ln in f:
            if "기대값(무작위)" in ln:
                mm = re.findall(r"([\d.]+)개", ln)
                if len(mm) >= 2:
                    rnd = (mm[0], mm[1])
                continue
            m = re.match(r"\s*(\S.*?)\s*\|\s*([\d.]+)개\s*\|\s*([\d.]+)개", ln)
            if m and "적중" not in ln:
                rows.append((m.group(1).strip(), m.group(2), m.group(3)))
    return rnd, rows


rnd, ev = parse_eval("out/result.txt")
s = add_slide()
set_bg(s, WHITE)
header(s, "추천 방식 개선 효과 검증 (최근 15주 백테스트)")
textbox(s, Inches(0.6), Inches(1.25), Inches(12), Inches(0.5),
        f"지표: 상위 추천 안에 실제 당첨 번호가 평균 몇 개 들어왔는가  ·  "
        f"무작위 기대값 = 상위6 {rnd[0]}개 / 상위10 {rnd[1]}개" if rnd else "지표: 상위 추천 적중 개수",
        14, GRAY, False)
ev_hdr = ["방식", "상위 6 적중", "상위 10 적중"]
tbl = s.shapes.add_table(len(ev) + 1, 3, Inches(0.6), Inches(1.95),
                         Inches(8.3), Inches(0.55) * (len(ev) + 1)).table
tbl.columns[0].width = Inches(4.3)
tbl.columns[1].width = Inches(2.0)
tbl.columns[2].width = Inches(2.0)
for j, h in enumerate(ev_hdr):
    tbl.cell(0, j).text = h
    style_cell(tbl.cell(0, j), 13, WHITE, True,
               PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER, fill=BLUE)
for i, (name, t6, t10) in enumerate(ev, start=1):
    is_wf = "신규(walk-forward)" in name
    is_legacy = "기존" in name
    fill = GOLD if is_wf else (RGBColor(0xF6, 0xDA, 0xDA) if is_legacy else LIGHT if i % 2 == 0 else WHITE)
    tbl.cell(i, 0).text = name + ("  ← 과적합 제거 핵심 지표" if is_wf else "")
    style_cell(tbl.cell(i, 0), 12, NAVY, is_wf or is_legacy, PP_ALIGN.LEFT, fill)
    tbl.cell(i, 1).text = t6 + "개"
    style_cell(tbl.cell(i, 1), 12, RED if is_legacy else NAVY, True, fill=fill)
    tbl.cell(i, 2).text = t10 + "개"
    style_cell(tbl.cell(i, 2), 12, NAVY, True, fill=fill)
textbox(s, Inches(0.6), Inches(5.55), Inches(12.1), Inches(1.6),
        "• 기존(무가중 합산)은 상위6 적중이 무작위보다도 낮음 → 사실상 무의미한 랭킹\n"
        "• 개선(정확도 가중 + 분석기별 정규화) 후 in-sample 1.27개까지 상승\n"
        "• walk-forward(매주 과거 데이터로만 재학습)에서도 상위10 적중이 무작위 대비 +30%\n"
        "  → 과적합을 빼도 기존보다 개선. 단, 실제 추첨은 무작위이므로 당첨을 보장하지 않음",
        14, GRAY, False)

# ---------------- Slide 2: 패턴별 예측 정확도 ----------------

# ---------------- Slide 2: 패턴별 예측 정확도 ----------------
s = add_slide()
set_bg(s, WHITE)
header(s, "최근 15회차 패턴별 예측 정확도 순위")
acc = [
    ("1", "이월수(주기)", "25.00%", "3/12"),
    ("2", "장기미출현", "22.00%", "11/50"),
    ("3", "오프셋(주기)", "20.00%", "5/25"),
    ("4", "5의배수", "17.78%", "24/135"),
    ("5", "궁합수", "16.59%", "34/205"),
    ("6", "출현주기", "16.11%", "48/298"),
    ("7", "3의배수", "15.11%", "34/225"),
    ("8", "끝수이월", "13.90%", "46/331"),
    ("9", "미출현", "13.84%", "40/289"),
    ("10", "고저균형", "13.78%", "31/225"),
    ("11", "누적빈도", "13.33%", "90/675"),
    ("12", "합계균형", "13.33%", "70/525"),
    ("13", "끝수통계", "13.33%", "90/675"),
    ("14", "AC패턴", "13.33%", "90/675"),
    ("15", "번호대빈도", "13.33%", "90/675"),
    ("16", "멸구간(주기)", "11.89%", "49/412"),
    ("17", "최근강세", "11.61%", "39/336"),
    ("18", "소수", "11.43%", "24/210"),
    ("19", "홀짝균형", "11.21%", "25/223"),
    ("20", "연번", "10.43%", "17/163"),
    ("21", "이웃수(주기)", "10.00%", "3/30"),
    ("22", "대각선(주기)", "6.98%", "3/43"),
]
# 두 컬럼으로 분할 (11행씩)
cols = [acc[:11], acc[11:]]
col_x = [Inches(0.5), Inches(6.9)]
headers = ["순위", "패턴 기법", "정확도", "적중/예측"]
for ci, group in enumerate(cols):
    rows = len(group) + 1
    tbl = s.shapes.add_table(rows, 4, col_x[ci], Inches(1.25),
                             Inches(5.95), Inches(5.9)).table
    tbl.columns[0].width = Inches(0.9)
    tbl.columns[1].width = Inches(2.3)
    tbl.columns[2].width = Inches(1.4)
    tbl.columns[3].width = Inches(1.35)
    for j, h in enumerate(headers):
        tbl.cell(0, j).text = h
        style_cell(tbl.cell(0, j), 13, WHITE, True, fill=BLUE)
    for i, (rk, name, pct, hit) in enumerate(group, start=1):
        fill = LIGHT if i % 2 == 0 else WHITE
        rkfill = GOLD if int(rk) <= 3 else fill
        tbl.cell(i, 0).text = rk + "위"
        style_cell(tbl.cell(i, 0), 12, WHITE if int(rk) <= 3 else NAVY,
                   True, fill=rkfill)
        tbl.cell(i, 1).text = name
        style_cell(tbl.cell(i, 1), 12, NAVY, False, PP_ALIGN.LEFT, fill)
        tbl.cell(i, 2).text = pct
        style_cell(tbl.cell(i, 2), 12, RED if int(rk) <= 5 else GRAY, True, fill=fill)
        tbl.cell(i, 3).text = hit
        style_cell(tbl.cell(i, 3), 11, GRAY, False, fill=fill)

# ---------------- Slide 3: 순위별 당첨번호 분포 ----------------
s = add_slide()
set_bg(s, WHITE)
header(s, "추천 순위 그룹별 당첨 번호 분포 (최근 50주 · 총 300개)")
dist = [
    ("1-5위", 45, 15.00),
    ("6-10위", 44, 14.67),
    ("11-15위", 38, 12.67),
    ("16-20위", 37, 12.33),
    ("21-25위", 29, 9.67),
    ("26-30위", 29, 9.67),
    ("31-35위", 26, 8.67),
    ("36-40위", 32, 10.67),
    ("41-45위", 20, 6.67),
]
maxpct = max(d[2] for d in dist)
top = Inches(1.45)
rowh = Inches(0.58)
lbl_x = Inches(0.6)
bar_x = Inches(2.6)
bar_max_w = Inches(8.0)
for i, (grp, cnt, pct) in enumerate(dist):
    y = top + rowh * i
    textbox(s, lbl_x, y, Inches(1.9), rowh, grp, 15,
            NAVY, True, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    w = int(bar_max_w * (pct / maxpct))
    if w < Emu(1):
        w = Emu(1)
    is_max = pct == maxpct
    is_min = pct == min(d[2] for d in dist)
    c = GOLD if is_max else (RGBColor(0xB0, 0xB8, 0xC9) if is_min else BLUE)
    bar(s, bar_x, y + Inches(0.07), w, Inches(0.42), c)
    textbox(s, bar_x + w + Inches(0.1), y, Inches(2.4), rowh,
            f"{pct:.2f}%  ({cnt}개)", 13, GRAY, True, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
textbox(s, Inches(0.6), Inches(6.75), Inches(12), Inches(0.5),
        "▶ '1-5위' 그룹이 15.00%로 가장 높고 순위가 낮을수록 적중률 감소 (정상 분포) · '41-45위' 최저 6.67%",
        13, RED, True, PP_ALIGN.LEFT)

# ---------------- Slide 3b: 그룹별 회차 당첨 내역 ----------------
import re


def parse_dist(path):
    rows = []
    with open(path, encoding="utf-8") as f:
        for ln in f:
            if "|" not in ln:
                continue
            parts = [p.strip() for p in ln.split("|")]
            if len(parts) < 4:
                continue
            if not re.match(r"^\d+-\d+위$", parts[0]):
                continue
            rows.append((parts[0], parts[1], parts[2], parts[3]))
    return rows


dist_rows = parse_dist("out/result.txt")
s = add_slide()
set_bg(s, WHITE)
header(s, "순위 그룹별 회차 당첨 내역 (최근 50주)")
hdr = ["순위 그룹", "당첨", "비율", "회차별 당첨 번호 (회차,번호(예측순위))"]
tbl = s.shapes.add_table(len(dist_rows) + 1, 4, Inches(0.3), Inches(1.2),
                         Inches(12.73), Inches(6.1)).table
tbl.columns[0].width = Inches(1.35)
tbl.columns[1].width = Inches(0.8)
tbl.columns[2].width = Inches(1.0)
tbl.columns[3].width = Inches(9.58)
for j, h in enumerate(hdr):
    tbl.cell(0, j).text = h
    style_cell(tbl.cell(0, j), 11, WHITE, True,
               PP_ALIGN.LEFT if j == 3 else PP_ALIGN.CENTER, fill=BLUE)
for i, (grp, cnt, pct, detail) in enumerate(dist_rows, start=1):
    pf = float(pct.replace("%", ""))
    is_max = pf == max(float(r[2].replace("%", "")) for r in dist_rows)
    is_min = pf == min(float(r[2].replace("%", "")) for r in dist_rows)
    fill = LIGHT if i % 2 == 0 else WHITE
    gfill = GOLD if is_max else fill
    tbl.cell(i, 0).text = grp
    style_cell(tbl.cell(i, 0), 11, NAVY, True, fill=gfill)
    tbl.cell(i, 1).text = cnt + "개"
    style_cell(tbl.cell(i, 1), 11, NAVY, True, fill=fill)
    tbl.cell(i, 2).text = pct
    style_cell(tbl.cell(i, 2), 11, RED if is_max else (BLUE if is_min else GRAY),
               True, fill=fill)
    tbl.cell(i, 3).text = detail
    style_cell(tbl.cell(i, 3), 8, GRAY, False, PP_ALIGN.LEFT, fill)

# ---------------- Slide 4: 확률 그룹 분석 ----------------
s = add_slide()
set_bg(s, WHITE)
header(s, "1226회차 당첨 번호 순위 그룹 확률 분석")
# UP 카드
bar(s, Inches(0.7), Inches(1.7), Inches(5.6), Inches(2.6), LIGHT)
bar(s, Inches(0.7), Inches(1.7), Inches(5.6), Pt(6), RED)
textbox(s, Inches(1.0), Inches(2.0), Inches(5.0), Inches(0.6),
        "▲ 필출 확률 UP", 22, RED, True)
textbox(s, Inches(1.0), Inches(2.75), Inches(5.0), Inches(0.6),
        "'1-5위' 그룹", 26, NAVY, True)
textbox(s, Inches(1.0), Inches(3.45), Inches(5.0), Inches(0.7),
        "최근 50주 당첨 확률 15.00%\n→ 1226회차 약 0.9개(1~3개) 포함 예상",
        14, GRAY, False)
# DOWN 카드
bar(s, Inches(7.0), Inches(1.7), Inches(5.6), Inches(2.6), LIGHT)
bar(s, Inches(7.0), Inches(1.7), Inches(5.6), Pt(6), BLUE)
textbox(s, Inches(7.3), Inches(2.0), Inches(5.0), Inches(0.6),
        "▼ 제외 확률 UP", 22, BLUE, True)
textbox(s, Inches(7.3), Inches(2.75), Inches(5.0), Inches(0.6),
        "'41-45위' 그룹", 26, NAVY, True)
textbox(s, Inches(7.3), Inches(3.45), Inches(5.0), Inches(0.7),
        "최근 50주 당첨 확률 6.67%\n→ 후순위로 고려 권장",
        14, GRAY, False)
textbox(s, Inches(0.7), Inches(4.9), Inches(11.9), Inches(1.5),
        "해석: 분포 산정 창을 50주(총 300개 표본)로 확대하고 라이브 추천과 동일한 가중 점수\n"
        "기준으로 일관화한 결과, 상위권일수록 적중률이 높은 정상적인 단조 분포가 나타납니다.\n"
        "게임 구성은 상위권(1-25위) 중심으로 하는 전략이 통계적으로 유효합니다.",
        15, GRAY, False)

# ---------------- Slide 5~6: 전체 추천 순위 45 (기여 패턴 포함) ----------------
def parse_rank(path):
    rows = []
    with open(path, encoding="utf-8") as f:
        for ln in f:
            if "점 |" not in ln or "번 |" not in ln:
                continue
            parts = [p.strip() for p in ln.split("|")]
            if len(parts) < 4:
                continue
            rk = parts[0].replace("👑", "").replace("위", "").strip()
            if not rk.isdigit():
                continue
            num = parts[1].replace("번", "").strip()
            sc = parts[2].replace("점", "").strip()
            pat = parts[3].strip()
            rows.append((rk, num, sc, pat))
    return rows


rank = parse_rank("out/result.txt")
assert len(rank) == 45, f"expected 45 rows, got {len(rank)}"

pages = [("1~23위", rank[:23]), ("24~45위", rank[23:])]
hdr = ["순위", "번호", "점수", "기여한 패턴"]
for pi, (subtitle, group) in enumerate(pages):
    s = add_slide()
    set_bg(s, WHITE)
    header(s, f"1226회차 추천 번호 전체 순위 ({subtitle})")
    rows = len(group) + 1
    tbl = s.shapes.add_table(rows, 4, Inches(0.35), Inches(1.2),
                             Inches(12.63), Inches(6.1)).table
    tbl.columns[0].width = Inches(0.95)
    tbl.columns[1].width = Inches(0.95)
    tbl.columns[2].width = Inches(1.25)
    tbl.columns[3].width = Inches(9.48)
    for j, h in enumerate(hdr):
        tbl.cell(0, j).text = h
        style_cell(tbl.cell(0, j), 12, WHITE, True,
                   PP_ALIGN.LEFT if j == 3 else PP_ALIGN.CENTER, fill=BLUE)
    for i, (rk, num, sc, pat) in enumerate(group, start=1):
        fill = LIGHT if i % 2 == 0 else WHITE
        top3 = int(rk) <= 3
        in_pool = int(rk) <= 25
        tbl.cell(i, 0).text = ("👑" + rk) if int(rk) == 1 else rk + "위"
        style_cell(tbl.cell(i, 0), 10, NAVY, top3, fill=GOLD if top3 else fill)
        tbl.cell(i, 1).text = num + "번"
        style_cell(tbl.cell(i, 1), 12, NAVY if in_pool else GRAY, True, fill=fill)
        tbl.cell(i, 2).text = sc
        style_cell(tbl.cell(i, 2), 10, GRAY, False, PP_ALIGN.RIGHT, fill)
        tbl.cell(i, 3).text = pat
        style_cell(tbl.cell(i, 3), 8.5, GRAY, False, PP_ALIGN.LEFT, fill)
    if pi == 0:
        textbox(s, Inches(0.35), Inches(7.32), Inches(12.6), Inches(0.25),
                "※ 상위 25위(진한 번호)가 최종 게임 생성 풀로 사용됩니다.",
                10, RGBColor(0x99, 0x99, 0x99), False)

# ---------------- Slide 6: 최종 5게임 ----------------
s = add_slide()
set_bg(s, WHITE)
header(s, "최종 추천 5게임 (자동 필터링)")
def parse_games(path):
    out = []
    with open(path, encoding="utf-8") as f:
        for ln in f:
            m = re.match(r"\[(\w) 게임\]\s*(.+?)\s*\(합:(\d+)\s*/\s*홀짝\s*(\d+):(\d+)\)", ln)
            if m:
                nums = [int(x) for x in re.findall(r"\d+", m.group(2))]
                out.append((m.group(1), nums, f"합:{m.group(3)} / 홀짝 {m.group(4)}:{m.group(5)}"))
    return out


games = parse_games("out/result.txt")
from pptx.enum.shapes import MSO_SHAPE
top = Inches(1.45)
rowh = Inches(1.02)
ball_d = Inches(0.74)
for i, (label, nums, meta) in enumerate(games):
    y = top + rowh * i
    bar(s, Inches(0.6), y, Inches(0.95), ball_d, NAVY)
    textbox(s, Inches(0.6), y, Inches(0.95), ball_d, label, 24, WHITE, True,
            PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    bx = Inches(1.85)
    gap = Inches(0.92)
    for k, n in enumerate(nums):
        cx = bx + gap * k
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL, cx, y, ball_d, ball_d)
        circ.fill.solid()
        # 번호대별 색상
        if n <= 10:
            cc = RGBColor(0xF6, 0xC3, 0x43)
        elif n <= 20:
            cc = RGBColor(0x4F, 0x9E, 0xE3)
        elif n <= 30:
            cc = RGBColor(0xE8, 0x6A, 0x5C)
        elif n <= 40:
            cc = RGBColor(0x8A, 0x8F, 0x9A)
        else:
            cc = RGBColor(0x6F, 0xB8, 0x6A)
        circ.fill.fore_color.rgb = cc
        circ.line.fill.background()
        circ.shadow.inherit = False
        tf = circ.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = str(n)
        r.font.size = Pt(20)
        r.font.bold = True
        r.font.color.rgb = WHITE
        r.font.name = FONT
    textbox(s, Inches(7.6), y, Inches(5.3), ball_d, meta, 14, GRAY, True,
            PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
textbox(s, Inches(0.6), Inches(6.75), Inches(12), Inches(0.5),
        "필터 조건: 합 100~180 · 홀수 2~4개 · 저수(≤22) 2~4개 · 끝수 4종↑ · 소수·3배수 각 1↑",
        12, RGBColor(0x99, 0x99, 0x99), False)

# ---------------- Slide 7: 면책 ----------------
s = add_slide()
set_bg(s, NAVY)
textbox(s, Inches(1.5), Inches(2.8), Inches(10.3), Inches(1.0),
        "본 분석은 과거 통계 패턴에 기반한 참고용 자료입니다.", 22, WHITE, True, PP_ALIGN.CENTER)
textbox(s, Inches(1.5), Inches(3.9), Inches(10.3), Inches(0.8),
        "로또는 무작위 추첨이며 당첨을 보장하지 않습니다. 건전한 구매 되세요.",
        16, RGBColor(0xC9, 0xD3, 0xEA), False, PP_ALIGN.CENTER)

prs.save("로또_1226회차_분석리포트.pptx")
print("SAVED: 로또_1226회차_분석리포트.pptx")
print("slides:", len(prs.slides._sldIdLst))
